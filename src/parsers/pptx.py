"""PPTX parser that extracts slide text, tables, and speaker notes."""

from __future__ import annotations

from pathlib import Path

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover - optional dependency
    Presentation = None

from .base import BaseParser, ParserError


class PptxParser(BaseParser):
    """Parse .pptx presentations with slide-level metadata."""

    source_type = "pptx"
    supported_extensions = (".pptx",)

    def parse(self, file_path: str | Path):
        path = Path(file_path)
        if not path.exists():
            raise ParserError(f"PPTX file not found: {path}")
        if Presentation is None:
            raise ParserError("python-pptx is required to parse .pptx files")

        try:
            presentation = Presentation(str(path))
        except Exception as exc:  # pragma: no cover - backend exceptions vary
            raise ParserError(f"Unable to open PPTX file {path}: {exc}") from exc

        output_lines: list[str] = []
        slides_metadata: list[dict[str, object]] = []

        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_lines = self._extract_slide_lines(slide)
            if not slide_lines:
                continue

            output_lines.append(f"Slide {slide_index}:")
            output_lines.extend(slide_lines)
            output_lines.append("")

            slides_metadata.append(
                {
                    "slide_number": slide_index,
                    "line_count": len(slide_lines),
                    "has_notes": any(line.startswith("Notes:") for line in slide_lines),
                }
            )

        output_lines = self._strip_outer_whitespace(output_lines)
        if not output_lines:
            raise ParserError(f"PPTX file has no readable content: {path}")

        metadata = {
            "format": self.source_type,
            "slide_count": len(slides_metadata),
            "slides": slides_metadata,
        }
        return self._build_document(
            text="\n".join(output_lines),
            file_path=path,
            source_type=self.source_type,
            metadata=metadata,
        )

    def _extract_slide_lines(self, slide) -> list[str]:
        lines: list[str] = []

        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                lines.append(text.strip())

            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    row_values = [
                        cell.text_frame.text.strip()
                        for cell in row.cells
                        if cell.text_frame.text.strip()
                    ]
                    if row_values:
                        lines.append(" | ".join(row_values))

        if getattr(slide, "has_notes_slide", False):
            notes_frame = slide.notes_slide.notes_text_frame
            notes_text = notes_frame.text.strip() if notes_frame else ""
            if notes_text:
                lines.append(f"Notes: {notes_text}")

        return [line for line in lines if line.strip()]
